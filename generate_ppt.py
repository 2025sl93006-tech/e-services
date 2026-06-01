from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x0D, 0x6E, 0xFD)   # Bootstrap primary
DARK_BLUE  = RGBColor(0x09, 0x4D, 0xB5)
DARK       = RGBColor(0x1E, 0x1E, 0x2E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
ACCENT     = RGBColor(0x0D, 0xCA, 0xF0)   # cyan
GREEN      = RGBColor(0x19, 0x87, 0x54)
ORANGE     = RGBColor(0xFD, 0x7E, 0x14)
GRAY       = RGBColor(0x6C, 0x75, 0x7D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Blue top bar with title + optional subtitle."""
    add_rect(slide, 0, 0, 13.33, 1.25, fill=BLUE)
    add_text(slide, title, 0.4, 0.12, 12.5, 0.65,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.78, 12.5, 0.38,
                 size=14, color=RGBColor(0xCC, 0xE5, 0xFF), italic=True)

def accent_bar(slide):
    """Thin cyan bottom accent line."""
    add_rect(slide, 0, 7.3, 13.33, 0.2, fill=ACCENT)

def section_badge(slide, label, l, t, color=BLUE):
    box = add_rect(slide, l, t, 2.2, 0.32, fill=color)
    add_text(slide, label, l + 0.08, t + 0.03, 2.1, 0.28,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, title, lines, title_color=BLUE,
         icon="", bg=LIGHT_GRAY):
    add_rect(slide, l, t, w, h, fill=bg, line=RGBColor(0xDE, 0xE2, 0xE6), line_w=Pt(1))
    txb = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.12), Inches(w-0.3), Inches(0.32))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = (icon + "  " if icon else "") + title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = title_color
    for line in lines:
        add_para(tf, line, size=11, color=DARK, space_before=Pt(3))

def bullet_box(slide, l, t, w, h, items, size=13, color=DARK,
               bullet="▸", bold_first=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(sl, 0, 0, 13.33, 0.18, fill=BLUE)
add_rect(sl, 0, 7.32, 13.33, 0.18, fill=ACCENT)

# decorative circles
for (cx, cy, sz, alpha) in [(10.5,1.2,2.5,DARK_BLUE),(11.8,5.5,3.2,RGBColor(0x12,0x52,0xBC))]:
    c = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = alpha
    c.line.fill.background()

add_text(sl, "e-Services", 1.5, 1.6, 10, 1.2, size=54, bold=True, color=WHITE)
add_text(sl, "Professional Home & Business Service Marketplace", 1.5, 2.85,
         10, 0.55, size=20, color=ACCENT, italic=True)

add_rect(sl, 1.5, 3.6, 4.5, 0.05, fill=BLUE)

for (label, val, yy) in [
    ("Platform",     "MERN Stack  (MongoDB · Express · React · Node.js)", 3.85),
    ("Architecture", "MVC  ·  REST API  ·  JWT Authentication",           4.3),
    ("Database",     "MongoDB  —  5 collections",                         4.75),
    ("Team",         "FSAD Course Project  —  2026",                      5.2),
]:
    add_text(sl, label + ":", 1.5, yy, 2.2, 0.4, size=13, bold=True,
             color=ACCENT)
    add_text(sl, val,         3.8, yy, 8.0, 0.4, size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Solution
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Problem & Solution",
             "Why e-Services was built")
accent_bar(sl)

# Problem box
add_rect(sl, 0.4, 1.45, 5.8, 5.4, fill=WHITE,
         line=RGBColor(0xDE, 0xE2, 0xE6), line_w=Pt(1))
add_rect(sl, 0.4, 1.45, 5.8, 0.42, fill=RGBColor(0xDC,0x35,0x45))
add_text(sl, "⚠  The Problem", 0.55, 1.49, 5.5, 0.36,
         size=14, bold=True, color=WHITE)

txb = sl.shapes.add_textbox(Inches(0.6), Inches(2.02), Inches(5.4), Inches(4.5))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for item in [
    "Finding reliable home service providers (plumbers, electricians, carpenters) is fragmented and unreliable.",
    "No transparent pricing — customers don't know the cost until the job is done.",
    "No way to verify provider quality before hiring.",
    "Booking is done via phone calls with no digital trail.",
    "No feedback mechanism to hold providers accountable.",
]:
    add_para(tf, f"✗  {item}", size=12, color=RGBColor(0x6C,0x75,0x7D), space_before=Pt(8))

# Solution box
add_rect(sl, 6.8, 1.45, 6.1, 5.4, fill=WHITE,
         line=RGBColor(0xDE, 0xE2, 0xE6), line_w=Pt(1))
add_rect(sl, 6.8, 1.45, 6.1, 0.42, fill=GREEN)
add_text(sl, "✔  Our Solution", 6.95, 1.49, 5.8, 0.36,
         size=14, bold=True, color=WHITE)

txb2 = sl.shapes.add_textbox(Inches(7.0), Inches(2.02), Inches(5.7), Inches(4.5))
txb2.word_wrap = True
tf2 = txb2.text_frame; tf2.word_wrap = True
for item in [
    "Centralised marketplace listing verified providers across 8 service categories.",
    "Transparent hourly rates and minimum charges shown upfront.",
    "Real star ratings (1–5) and written reviews from verified customers.",
    "Digital booking with auto-generated order reference and email confirmation.",
    "Admin panel to manage providers, categories, and order lifecycle.",
]:
    add_para(tf2, f"✔  {item}", size=12, color=DARK, space_before=Pt(8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Features
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Key Features", "What the platform offers")
accent_bar(sl)

features = [
    ("🔐", "Authentication",        BLUE,   ["JWT-based register & login", "Role-based access (admin / enduser)", "Password change & profile update"]),
    ("🔍", "Smart Search & Filter",  DARK_BLUE, ["Filter by category, location, cost range", "Full-text keyword search", "Availability status filter"]),
    ("📋", "Service Booking",        GREEN,  ["Live cost preview before confirming", "Auto order ref  ES-YYYY-XXXX", "Email confirmation via Nodemailer"]),
    ("⭐", "Ratings & Reviews",      ORANGE, ["1–5 star interactive rating", "One review per user per provider", "Auto-updates provider avg rating"]),
    ("🛡️", "Admin Panel",            RGBColor(0x6F,0x42,0xC1), ["CRUD for categories & providers", "Paginated order management", "Status workflow: pending → confirmed → in-progress → completed"]),
    ("📊", "Dashboards",             RGBColor(0x0D,0xCA,0xF0), ["User: KPI cards + recent orders", "Admin: platform-wide stats", "Cancel orders with reason"]),
]

cols = [(0.3, 1.4), (4.55, 1.4), (8.8, 1.4),
        (0.3, 4.1), (4.55, 4.1), (8.8, 4.1)]
for i, (icon, title, col, lines) in enumerate(features):
    lx, ty = cols[i]
    add_rect(sl, lx, ty, 3.9, 2.5, fill=WHITE,
             line=RGBColor(0xDE,0xE2,0xE6), line_w=Pt(1))
    add_rect(sl, lx, ty, 3.9, 0.45, fill=col)
    add_text(sl, f"{icon}  {title}", lx+0.12, ty+0.06, 3.7, 0.35,
             size=13, bold=True, color=WHITE)
    txb = sl.shapes.add_textbox(Inches(lx+0.15), Inches(ty+0.55),
                                 Inches(3.65), Inches(1.85))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for j, ln in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {ln}"
        run.font.size = Pt(11)
        run.font.color.rgb = DARK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
slide_header(sl, "Technology Stack", "Libraries, frameworks and tools used")
add_rect(sl, 0, 7.3, 13.33, 0.2, fill=ACCENT)

layers = [
    ("FRONTEND",   BLUE,   [
        ("React 18",            "UI library with hooks"),
        ("React Router v6",     "Client-side routing"),
        ("Bootstrap 5",         "Responsive styling"),
        ("Formik + Yup",        "Form management & validation"),
        ("Axios",               "HTTP client with interceptors"),
        ("React-Toastify",      "Toast notifications"),
    ]),
    ("BACKEND",    GREEN,  [
        ("Node.js + Express",   "REST API server  (port 5001)"),
        ("Mongoose",            "MongoDB ODM"),
        ("bcryptjs",            "Password hashing"),
        ("jsonwebtoken",        "JWT auth tokens"),
        ("Nodemailer",          "Order confirmation emails"),
        ("swagger-jsdoc/ui",    "Auto-generated API docs"),
    ]),
    ("DATABASE &\nDEVOPS", ORANGE, [
        ("MongoDB",             "Primary database"),
        ("5 Collections",       "User · Category · ServiceProvider · Order · Review"),
        ("Compound Index",      "{user, provider} on Review — unique"),
        ("Text Index",          "name + description + location on Provider"),
        ("Soft Deletes",        "isActive / isVisible flags throughout"),
        ("Seed Script",         "npm run seed — 3 users, 8 cats, 16 providers"),
    ]),
]

for i, (label, col, items) in enumerate(layers):
    lx = 0.3 + i * 4.35
    add_rect(sl, lx, 1.35, 4.0, 0.45, fill=col)
    add_text(sl, label, lx+0.1, 1.38, 3.8, 0.4,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (tech, desc) in enumerate(items):
        ty = 1.95 + j * 0.82
        add_rect(sl, lx, ty, 4.0, 0.72,
                 fill=RGBColor(0x2A,0x2A,0x3E),
                 line=RGBColor(0x3A,0x3A,0x5E), line_w=Pt(1))
        add_text(sl, tech, lx+0.15, ty+0.06, 3.7, 0.3,
                 size=12, bold=True, color=col)
        add_text(sl, desc, lx+0.15, ty+0.36, 3.7, 0.28,
                 size=10, color=RGBColor(0xAA,0xAA,0xCC))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "System Architecture", "MVC pattern — Request → Route → Controller → Model → Response")
accent_bar(sl)

boxes = [
    (0.25, 2.2, 2.2, 3.8, BLUE,   "CLIENT\n(React 18)",
     ["HomePage", "ServicesPage", "ServiceDetailPage", "OrderServicePage",
      "UserDashboard", "MyOrdersPage", "AdminDashboard", "Manage Pages"]),
    (3.2,  2.2, 2.4, 3.8, DARK_BLUE, "API LAYER\n(Axios)",
     ["axiosInstance.js", "Auto-attach JWT", "401 → logout", "authApi.js",
      "categoryApi.js", "serviceApi.js", "orderApi.js", "reviewApi.js"]),
    (6.35, 2.2, 2.4, 3.8, GREEN, "EXPRESS SERVER\n(MVC)",
     ["Routes (5 files)", "Controllers (5)", "authMiddleware", "requireAdmin",
      "emailService", "swagger-ui", "errorHandler", "port 5001"]),
    (9.45, 2.2, 3.55, 3.8, RGBColor(0xFD,0x7E,0x14), "MONGODB\n(5 Collections)",
     ["User", "Category", "ServiceProvider", "Order  (orderRef ES-YYYY-XXXX)",
      "Review  (compound idx)", "averageRating denorm.", "Soft deletes", "Text indexes"]),
]

for (lx, ty, w, h, col, title, items) in boxes:
    add_rect(sl, lx, ty, w, h, fill=WHITE,
             line=col, line_w=Pt(2))
    add_rect(sl, lx, ty, w, 0.55, fill=col)
    add_text(sl, title, lx+0.1, ty+0.06, w-0.2, 0.46,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(Inches(lx+0.12), Inches(ty+0.65),
                                 Inches(w-0.24), Inches(h-0.75))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for k, item in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(10)
        run.font.color.rgb = DARK

# arrows between boxes
arrow_y = 4.1
for ax in [2.5, 5.65, 8.8]:
    add_rect(sl, ax, arrow_y-0.05, 0.65, 0.1, fill=GRAY)
    add_text(sl, "▶", ax+0.48, arrow_y-0.2, 0.3, 0.35, size=12, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Database Design
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Database Design", "MongoDB collections and relationships")
accent_bar(sl)

models = [
    ("User", BLUE, [
        "name, email, password (bcrypt)",
        "role: admin | enduser",
        "phone, address, isActive",
        "comparePassword() method",
    ]),
    ("Category", GREEN, [
        "name (unique, required)",
        "description, imageUrl",
        "isActive (soft delete)",
    ]),
    ("ServiceProvider", ORANGE, [
        "name, category (ref), location",
        "chargePerHour, minCharge",
        "availability: available|busy|unavailable",
        "averageRating, totalReviews (denorm.)",
        "Text index on name+desc+location",
    ]),
    ("Order", RGBColor(0x6F,0x42,0xC1), [
        "orderRef: ES-YYYY-XXXX (pre-save hook)",
        "user, provider, category (refs)",
        "serviceDate, serviceAddress, estimatedHours",
        "estimatedCost = max(rate×hrs, minCharge)",
        "orderStatus: pending→confirmed→in-progress→completed|cancelled",
        "emailSent, cancellationReason",
    ]),
    ("Review", RGBColor(0xDC,0x35,0x45), [
        "user, provider, order (refs)",
        "rating 1–5, comment, isVisible",
        "Compound unique index {user, provider}",
        "recomputeRating() static method",
        "Aggregates avg → updates ServiceProvider",
    ]),
]

positions = [(0.25,1.4),(2.85,1.4),(5.45,1.4),(0.25,4.05),(6.5,4.05)]
widths    = [2.4, 2.4, 2.4, 5.9, 6.55]
heights   = [2.4, 2.4, 2.4, 2.9, 2.9]

for i, (name, col, fields) in enumerate(models):
    lx, ty = positions[i]
    w, h = widths[i], heights[i]
    add_rect(sl, lx, ty, w, h, fill=WHITE,
             line=col, line_w=Pt(2))
    add_rect(sl, lx, ty, w, 0.4, fill=col)
    add_text(sl, name, lx+0.1, ty+0.05, w-0.2, 0.32,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(Inches(lx+0.12), Inches(ty+0.48),
                                 Inches(w-0.24), Inches(h-0.55))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for k, f in enumerate(fields):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {f}"
        run.font.size = Pt(10 if w > 3 else 10)
        run.font.color.rgb = DARK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — API Endpoints
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
slide_header(sl, "REST API — 27 Endpoints",
             "Documented via Swagger UI at  http://localhost:5001/api-docs")
add_rect(sl, 0, 7.3, 13.33, 0.2, fill=ACCENT)

groups = [
    ("/auth",        BLUE,   ["POST  /register", "POST  /login", "GET   /me", "PUT   /me", "PUT   /change-password"]),
    ("/categories",  GREEN,  ["GET   /  (public)", "GET   /:id  (public)", "POST  /  (admin)", "PUT   /:id  (admin)", "DELETE /:id  (admin)"]),
    ("/providers",   ORANGE, ["GET   /  + filters", "GET   /category/:id", "GET   /:id", "POST  /  (admin)", "PUT   /:id  (admin)", "DELETE /:id"]),
    ("/orders",      RGBColor(0x6F,0x42,0xC1), ["POST  /  (create)", "GET   /  (admin paged)", "GET   /my", "GET   /:id", "PUT   /:id/cancel", "PUT   /:id/status (admin)"]),
    ("/reviews",     RGBColor(0xDC,0x35,0x45), ["POST  /  (create)", "GET   /  (admin)", "GET   /my", "GET   /provider/:id (public)", "PUT   /:id  (owner)", "DELETE /:id  (admin)"]),
]

for i, (prefix, col, endpoints) in enumerate(groups):
    lx = 0.25 + i * 2.62
    add_rect(sl, lx, 1.38, 2.45, 0.42, fill=col)
    add_text(sl, prefix, lx+0.08, 1.41, 2.3, 0.36,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, 1.8, 2.45, 5.3, fill=RGBColor(0x1E,0x1E,0x2E),
             line=col, line_w=Pt(1))
    txb = sl.shapes.add_textbox(Inches(lx+0.12), Inches(1.9),
                                 Inches(2.25), Inches(5.1))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for k, ep in enumerate(endpoints):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        parts = ep.split("  ", 1)
        method = parts[0].strip()
        rest   = parts[1] if len(parts) > 1 else ""
        run = p.add_run()
        run.text = method + "  "
        run.font.size = Pt(11)
        run.font.bold = True
        method_colors = {"POST": RGBColor(0x19,0x87,0x54), "GET": BLUE,
                         "PUT": ORANGE, "DELETE": RGBColor(0xDC,0x35,0x45)}
        run.font.color.rgb = method_colors.get(method, col)
        run2 = p.add_run()
        run2.text = rest
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(0xCC,0xCC,0xDD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — User Flow
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "User Flow", "End-to-end journey for a customer")
accent_bar(sl)

steps = [
    (BLUE,   "1",  "Register / Login",     ["Fill name, email, password", "JWT token issued on success", "Role = enduser"]),
    (GREEN,  "2",  "Browse Services",       ["View homepage categories", "Filter by location / cost", "See star ratings"]),
    (ORANGE, "3",  "View Provider Detail",  ["See provider profile", "Read customer reviews", "Check availability & pricing"]),
    (RGBColor(0x6F,0x42,0xC1), "4", "Place Order", ["Select date, address, hours", "Live cost preview shown", "Email confirmation sent"]),
    (RGBColor(0xDC,0x35,0x45), "5", "Track Order",  ["View order in My Orders", "Cancel if still pending", "Status updated by admin"]),
    (RGBColor(0x0D,0xCA,0xF0), "6", "Leave Review", ["Available after completed order", "1–5 star rating + comment", "Provider avg rating updates live"]),
]

for i, (col, num, title, lines) in enumerate(steps):
    lx = 0.3 + i * 2.15
    # circle
    circ = sl.shapes.add_shape(9, Inches(lx+0.7), Inches(1.38),
                                Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    add_text(sl, num, lx+0.73, 1.4, 0.56, 0.52,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # connector (except last)
    if i < 5:
        add_rect(sl, lx+1.35, 1.64, 0.82, 0.08, fill=GRAY)
        add_text(sl, "▶", lx+2.1, 1.55, 0.2, 0.3, size=9, color=GRAY)
    # card
    add_rect(sl, lx, 2.2, 2.0, 4.5, fill=WHITE,
             line=col, line_w=Pt(2))
    add_rect(sl, lx, 2.2, 2.0, 0.42, fill=col)
    add_text(sl, title, lx+0.08, 2.24, 1.88, 0.36,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb = sl.shapes.add_textbox(Inches(lx+0.12), Inches(2.72),
                                 Inches(1.8), Inches(3.8))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"▸  {ln}"
        run.font.size = Pt(11)
        run.font.color.rgb = DARK

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Admin Flow
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Admin Panel", "Platform management capabilities")
accent_bar(sl)

# Left: what admin can do
add_rect(sl, 0.3, 1.42, 6.0, 5.6, fill=WHITE,
         line=BLUE, line_w=Pt(2))
add_rect(sl, 0.3, 1.42, 6.0, 0.42, fill=BLUE)
add_text(sl, "Admin Capabilities", 0.45, 1.46, 5.8, 0.36,
         size=14, bold=True, color=WHITE)

txb = sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.6), Inches(4.8))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
admin_items = [
    ("Manage Categories",     "Create, edit, deactivate service categories (Civil, Electrical, etc.)"),
    ("Manage Providers",      "Add/edit/deactivate providers with full detail — category, rates, location, availability"),
    ("Manage Orders",         "View all orders paginated, filter by status, advance order through status workflow"),
    ("Moderate Reviews",      "View all reviews (including hidden), soft-delete inappropriate content"),
    ("Platform Dashboard",    "KPI cards: total orders, pending orders, provider count, category count"),
]
for label, desc in admin_items:
    add_para(tf, f"▸  {label}", size=13, bold=True, color=BLUE, space_before=Pt(10))
    add_para(tf, f"   {desc}",  size=11, color=DARK, space_before=Pt(2))

# Right: access control
add_rect(sl, 6.8, 1.42, 6.1, 2.6, fill=WHITE,
         line=GREEN, line_w=Pt(2))
add_rect(sl, 6.8, 1.42, 6.1, 0.42, fill=GREEN)
add_text(sl, "Access Control", 6.95, 1.46, 5.8, 0.36,
         size=14, bold=True, color=WHITE)
txb2 = sl.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.7), Inches(1.9))
txb2.word_wrap = True
tf2 = txb2.text_frame; tf2.word_wrap = True
for item in [
    "JWT payload carries role — no DB lookup per request",
    "verifyToken middleware validates Bearer token",
    "requireAdmin middleware checks role === 'admin'",
    "AdminRoute component blocks non-admins in React",
]:
    add_para(tf2, f"▸  {item}", size=11, color=DARK, space_before=Pt(6))

# Right: order status workflow
add_rect(sl, 6.8, 4.22, 6.1, 2.8, fill=WHITE,
         line=ORANGE, line_w=Pt(2))
add_rect(sl, 6.8, 4.22, 6.1, 0.42, fill=ORANGE)
add_text(sl, "Order Status Workflow  (Admin driven)", 6.95, 4.26, 5.8, 0.36,
         size=13, bold=True, color=WHITE)
statuses = [
    ("pending",     ORANGE),
    ("confirmed",   ACCENT),
    ("in-progress", BLUE),
    ("completed",   GREEN),
    ("cancelled",   RGBColor(0xDC,0x35,0x45)),
]
sx = 7.0
for j, (st, sc) in enumerate(statuses):
    add_rect(sl, sx, 4.82, 1.0, 0.32, fill=sc)
    add_text(sl, st, sx+0.04, 4.84, 0.95, 0.28,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if j < 4:
        add_text(sl, "→", sx+1.02, 4.82, 0.2, 0.32,
                 size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    sx += 1.2
add_text(sl, "User can cancel while status is  pending  or  confirmed",
         7.0, 5.32, 5.7, 0.32, size=11, italic=True, color=GRAY)
add_text(sl, "Admin advances status using  PUT /orders/:id/status",
         7.0, 5.68, 5.7, 0.32, size=11, italic=True, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Security & Technical Highlights
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
slide_header(sl, "Security & Technical Highlights",
             "Design decisions that make the system robust")
add_rect(sl, 0, 7.3, 13.33, 0.2, fill=ACCENT)

highlights = [
    (BLUE,   "🔐", "JWT Authentication",
     ["Token signed with HS256 + JWT_SECRET", "Payload: { id, role } — minimal, no PII", "Expiry configurable via JWT_EXPIRE (.env)", "Auto-logout on 401 via Axios interceptor"]),
    (GREEN,  "🛡️", "Password Security",
     ["bcryptjs with salt rounds = 10", "Passwords never returned in API responses", ".select('-password') on all user queries", "Change-password requires current password"]),
    (ORANGE, "📊", "Rating Integrity",
     ["Compound unique index prevents duplicate reviews", "recomputeRating() aggregation on every change", "averageRating stored denormalized on provider", "Rounded to 1 decimal place"]),
    (ACCENT, "📧", "Non-blocking Email",
     ["sendOrderConfirmation() called with .then().catch()", "Email failure never fails the order response", "emailSent flag tracked on Order document", "HTML table template with order details"]),
    (RGBColor(0x6F,0x42,0xC1), "🗄️", "Data Integrity",
     ["Soft deletes everywhere (isActive/isVisible)", "Mongoose validators + runValidators: true", "Server-side cost calculation — client not trusted", "Order ref generated in pre-save hook"]),
    (RGBColor(0xDC,0x35,0x45), "📝", "API Documentation",
     ["swagger-jsdoc reads @swagger JSDoc in routes/", "swagger-ui-express serves at /api-docs", "All schemas defined in swagger.js components", "persistAuthorization: true in Swagger config"]),
]

for i, (col, icon, title, points) in enumerate(highlights):
    lx = 0.25 + (i % 3) * 4.35
    ty = 1.42 + (i // 3) * 2.85
    add_rect(sl, lx, ty, 4.0, 2.65,
             fill=RGBColor(0x1E,0x1E,0x2E),
             line=col, line_w=Pt(2))
    add_rect(sl, lx, ty, 4.0, 0.44, fill=col)
    add_text(sl, f"{icon}  {title}", lx+0.1, ty+0.06, 3.8, 0.36,
             size=13, bold=True, color=WHITE)
    txb = sl.shapes.add_textbox(Inches(lx+0.15), Inches(ty+0.55),
                                 Inches(3.75), Inches(2.0))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for k, pt in enumerate(points):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▸  {pt}"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xCC,0xCC,0xDD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Project Structure
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Project Structure",
             "Directory layout of the full-stack application")
accent_bar(sl)

# Server tree
add_rect(sl, 0.3, 1.42, 5.9, 5.6, fill=WHITE,
         line=BLUE, line_w=Pt(2))
add_rect(sl, 0.3, 1.42, 5.9, 0.42, fill=BLUE)
add_text(sl, "server/", 0.45, 1.46, 5.6, 0.36,
         size=14, bold=True, color=WHITE)
server_tree = [
    ("server.js",         "Entry point — mounts routes + Swagger"),
    ("swagger.js",        "OpenAPI 3.0 spec definition"),
    ("seedData.js",       "DB seeder script"),
    ("config/",           "db.js, nodemailer.js"),
    ("middleware/",       "authMiddleware.js, errorHandler.js"),
    ("models/",           "User · Category · ServiceProvider · Order · Review"),
    ("controllers/",      "auth · category · provider · order · review"),
    ("routes/",           "5 route files with @swagger JSDoc"),
    ("services/",         "emailService.js"),
    ("utils/",            "generateToken.js"),
]
txb = sl.shapes.add_textbox(Inches(0.5), Inches(1.95), Inches(5.5), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for k, (f, d) in enumerate(server_tree):
    p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = f"  {f:<28}"
    run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = BLUE
    run2 = p.add_run(); run2.text = d
    run2.font.size = Pt(10); run2.font.color.rgb = GRAY

# Client tree
add_rect(sl, 6.8, 1.42, 6.1, 5.6, fill=WHITE,
         line=GREEN, line_w=Pt(2))
add_rect(sl, 6.8, 1.42, 6.1, 0.42, fill=GREEN)
add_text(sl, "client/src/", 6.95, 1.46, 5.8, 0.36,
         size=14, bold=True, color=WHITE)
client_tree = [
    ("App.jsx",              "Router + AuthProvider + ToastContainer"),
    ("index.js",             "ReactDOM.createRoot entry point"),
    ("index.css",            "Global styles + star/card CSS"),
    ("api/",                 "axiosInstance · authApi · categoryApi · serviceApi · orderApi · reviewApi"),
    ("context/AuthContext",  "user, token, login(), logout(), isAdmin()"),
    ("components/common/",   "PrivateRoute · AdminRoute · StarRating"),
    ("components/layout/",   "Navbar · Footer"),
    ("pages/ (13 pages)",    "Home · Services · Detail · Login · Register · UserDashboard · OrderService · MyOrders · AdminDashboard · ManageCategories · ManageProviders · ManageOrders"),
]
txb2 = sl.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.9))
txb2.word_wrap = True
tf2 = txb2.text_frame; tf2.word_wrap = True
for k, (f, d) in enumerate(client_tree):
    p = tf2.paragraphs[0] if k == 0 else tf2.add_paragraph()
    p.space_before = Pt(4)
    run = p.add_run(); run.text = f"  {f:<26}"
    run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = GREEN
    run2 = p.add_run(); run2.text = d
    run2.font.size = Pt(10); run2.font.color.rgb = GRAY

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Requirements Coverage
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
slide_header(sl, "Requirements Coverage",
             "All functional requirements satisfied")
accent_bar(sl)

reqs = [
    ("R1", "User Registration & Login",      BLUE,   "User model with bcrypt · JWT auth · register/login pages · AuthContext"),
    ("R2", "Filter by Category / Cost / Location", GREEN, "providerController.getAll with multi-filter: category, location (regex), minCost/maxCost, availability, $text search · ServicesPage filter panel"),
    ("R3", "Two User Roles",                 ORANGE, "User.role enum (admin | enduser) · requireAdmin middleware · AdminRoute component · role-based navbar dropdowns"),
    ("R4", "Admin Manage Providers & Categories", RGBColor(0x6F,0x42,0xC1), "Full CRUD on /providers and /categories · ManageProvidersPage · ManageCategoriesPage · Formik + Yup modals"),
    ("R5", "Search + Order Service",         RGBColor(0x0D,0xCA,0xF0), "$text index on ServiceProvider · ServicesPage search box · OrderServicePage · order ref generation · email confirmation"),
    ("R6", "Customer Feedback",              RGBColor(0xDC,0x35,0x45), "Review model · reviewController · ServiceDetailPage review form · MyOrdersPage inline review modal · isVisible soft delete"),
    ("R7", "Star Ratings",                   RGBColor(0xFD,0x7E,0x14), "Review.rating (1–5) · StarRating.jsx component (display + interactive) · recomputeRating() static method · averageRating on provider"),
]

for i, (ref, title, col, impl) in enumerate(reqs):
    ty = 1.42 + i * 0.84
    add_rect(sl, 0.3, ty, 12.7, 0.76, fill=WHITE,
             line=RGBColor(0xDE,0xE2,0xE6), line_w=Pt(1))
    add_rect(sl, 0.3, ty, 0.72, 0.76, fill=col)
    add_text(sl, ref, 0.31, ty+0.18, 0.7, 0.4,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 1.12, ty+0.06, 3.5, 0.32,
             size=12, bold=True, color=col)
    add_text(sl, impl, 1.12, ty+0.38, 11.6, 0.3,
             size=10, color=GRAY)
    # tick
    add_text(sl, "✔", 12.7, ty+0.18, 0.38, 0.36,
             size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Demo & How to Run
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
slide_header(sl, "Demo & How to Run",
             "Setup instructions and test accounts")
add_rect(sl, 0, 7.3, 13.33, 0.2, fill=ACCENT)

# Commands box
add_rect(sl, 0.3, 1.42, 6.0, 5.6, fill=RGBColor(0x1E,0x1E,0x2E),
         line=BLUE, line_w=Pt(2))
add_rect(sl, 0.3, 1.42, 6.0, 0.42, fill=BLUE)
add_text(sl, "⚙  Setup Commands", 0.45, 1.46, 5.8, 0.36,
         size=14, bold=True, color=WHITE)

cmds = [
    ("# 1. Start MongoDB", ""),
    ("mongod", ""),
    ("", ""),
    ("# 2. Start Server", ""),
    ("cd e-services/server", ""),
    ("npm install", "installs dependencies"),
    ("npm run seed", "wipes DB + loads test data"),
    ("npm run dev", "starts API on  :5001"),
    ("", ""),
    ("# 3. Start Client", ""),
    ("cd e-services/client", ""),
    ("npm install", ""),
    ("npm start", "starts React on  :3001"),
    ("", ""),
    ("# 4. API Documentation", ""),
    ("http://localhost:5001/api-docs", "Swagger UI"),
]
txb = sl.shapes.add_textbox(Inches(0.5), Inches(1.98), Inches(5.6), Inches(4.9))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
for k, (cmd, comment) in enumerate(cmds):
    p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    if not cmd:
        run = p.add_run(); run.text = " "; run.font.size = Pt(6)
        continue
    if cmd.startswith("#"):
        run = p.add_run(); run.text = cmd
        run.font.size = Pt(10); run.font.italic = True
        run.font.color.rgb = GRAY
    else:
        run = p.add_run(); run.text = cmd
        run.font.size = Pt(12); run.font.bold = True
        run.font.color.rgb = ACCENT
        if comment:
            run2 = p.add_run(); run2.text = f"   # {comment}"
            run2.font.size = Pt(10); run2.font.italic = True
            run2.font.color.rgb = GRAY

# Credentials box
add_rect(sl, 6.8, 1.42, 6.1, 2.7, fill=RGBColor(0x1E,0x1E,0x2E),
         line=GREEN, line_w=Pt(2))
add_rect(sl, 6.8, 1.42, 6.1, 0.42, fill=GREEN)
add_text(sl, "🔑  Test Accounts  (after npm run seed)", 6.95, 1.46, 5.8, 0.36,
         size=13, bold=True, color=WHITE)
creds = [
    ("Admin",  "admin@eservices.com",  "admin123", ORANGE),
    ("User 1", "rahul@example.com",    "user123",  BLUE),
    ("User 2", "priya@example.com",    "user123",  ACCENT),
]
for j, (role, email, pwd, cc) in enumerate(creds):
    ty2 = 2.06 + j * 0.72
    add_rect(sl, 7.0, ty2, 1.0, 0.52, fill=cc)
    add_text(sl, role, 7.02, ty2+0.1, 0.96, 0.34,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, email, 8.1, ty2+0.04, 2.9, 0.24, size=11, color=WHITE)
    add_text(sl, pwd,   8.1, ty2+0.28, 2.9, 0.22, size=11, color=GRAY, italic=True)

# URLs box
add_rect(sl, 6.8, 4.3, 6.1, 2.72, fill=RGBColor(0x1E,0x1E,0x2E),
         line=ACCENT, line_w=Pt(2))
add_rect(sl, 6.8, 4.3, 6.1, 0.42, fill=ACCENT)
add_text(sl, "🌐  Application URLs", 6.95, 4.34, 5.8, 0.36,
         size=13, bold=True, color=DARK)
urls = [
    ("React App",    "http://localhost:3001",           BLUE),
    ("API Server",   "http://localhost:5001/api",        GREEN),
    ("Swagger Docs", "http://localhost:5001/api-docs",   ORANGE),
    ("Health Check", "http://localhost:5001/api/health", GRAY),
]
txb3 = sl.shapes.add_textbox(Inches(7.0), Inches(4.88), Inches(5.7), Inches(2.0))
txb3.word_wrap = True
tf3 = txb3.text_frame; tf3.word_wrap = True
for k2, (label, url, uc) in enumerate(urls):
    p = tf3.paragraphs[0] if k2 == 0 else tf3.add_paragraph()
    p.space_before = Pt(7)
    run = p.add_run(); run.text = f"{label:<16}"
    run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = uc
    run2 = p.add_run(); run2.text = url
    run2.font.size = Pt(11); run2.font.color.rgb = WHITE

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(sl, 0, 0,     13.33, 0.18, fill=BLUE)
add_rect(sl, 0, 7.32,  13.33, 0.18, fill=ACCENT)

# decorative
for (cx, cy, sz, col) in [(9.5,0.5,3.5,DARK_BLUE),(10.5,4.5,4.0,RGBColor(0x0A,0x40,0x90))]:
    c = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = col
    c.line.fill.background()

add_text(sl, "Thank You", 1.5, 1.5, 9, 1.2,
         size=54, bold=True, color=WHITE)
add_text(sl, "e-Services — Professional Home & Business Service Marketplace",
         1.5, 2.85, 9, 0.55, size=18, color=ACCENT, italic=True)

add_rect(sl, 1.5, 3.6, 4.5, 0.06, fill=BLUE)

summary = [
    ("Stack",       "MERN  (MongoDB · Express · React 18 · Node.js)"),
    ("API",         "27 REST endpoints  ·  Swagger UI at /api-docs"),
    ("Collections", "5 MongoDB models with compound index & text search"),
    ("Pages",       "13 React pages  ·  2 roles  ·  JWT authentication"),
    ("Patterns",    "MVC  ·  Soft deletes  ·  Denormalised ratings  ·  Non-blocking email"),
]
for (label, val) in summary:
    yy = 3.85 + summary.index((label, val)) * 0.48
    add_text(sl, label + ":", 1.5, yy, 1.9, 0.4,
             size=13, bold=True, color=ACCENT)
    add_text(sl, val, 3.5, yy, 7.5, 0.4,
             size=13, color=WHITE)

add_text(sl, "Questions?", 1.5, 6.5, 5, 0.5,
         size=20, bold=True, color=BLUE)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/I578377/Desktop/Assignments/FSAD/FSAD Project/e-services/eServices_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
